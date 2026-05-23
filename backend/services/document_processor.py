import os
import re
from pathlib import Path
from typing import List, Tuple

import pdfplumber
from pypdf import PdfReader


def extract_text_from_file(filepath: str, file_type: str) -> Tuple[str, int]:
    """Extract text and page count from uploaded file."""
    ext = file_type.lower()
    text = ""
    page_count = 0

    try:
        if ext == ".pdf":
            text, page_count = _extract_pdf(filepath)
        elif ext in (".docx", ".doc"):
            text, page_count = _extract_docx(filepath)
        elif ext == ".pptx":
            text, page_count = _extract_pptx(filepath)
        elif ext == ".txt":
            text, page_count = _extract_txt(filepath)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        raise RuntimeError(f"Text extraction failed: {e}")

    text = clean_text(text)
    return text, page_count


def _extract_pdf(filepath: str) -> Tuple[str, int]:
    texts = []
    page_count = 0
    try:
        with pdfplumber.open(filepath) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
    except Exception:
        reader = PdfReader(filepath)
        page_count = len(reader.pages)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                texts.append(t)

    combined = "\n\n".join(texts)

    # Detect scanned PDFs: strip known watermarks and check remaining content
    content_check = combined
    for watermark in ["Downloaded from sajhanotes.com", "Downloaded from"]:
        content_check = content_check.replace(watermark, "")
    content_check = content_check.strip()

    # If fewer than 200 real characters were extracted, assume scanned image PDF
    if len(content_check) < 200 and page_count > 0:
        print(f"[INFO] Scanned PDF detected — only {len(content_check)} chars from text layer. Running AI OCR...")
        try:
            return _extract_pdf_with_ocr(filepath)
        except Exception as ocr_err:
            print(f"[WARN] AI OCR failed: {ocr_err}. Returning limited text.")

    return combined, page_count


def _extract_pdf_with_ocr(filepath: str) -> Tuple[str, int]:
    """
    Render each page of a scanned PDF as an image and extract text
    using Groq's vision model. No Tesseract required.
    """
    import fitz  # PyMuPDF
    from services.ai_client import generate_with_image

    doc = fitz.open(filepath)
    page_count = len(doc)
    texts = []

    OCR_PROMPT = (
        "You are an OCR tool. Output ONLY the raw text from this exam paper image. "
        "No explanations, no steps, no commentary — just the text exactly as printed. "
        "Include every question number, question text, all options (A/B/C/D), "
        "marks in brackets, group headers, and any instructions. "
        "Preserve the original structure and wording faithfully."
    )

    for i, page in enumerate(doc):
        # Render at 200 DPI in grayscale for smaller payload
        mat = fitz.Matrix(200 / 72, 200 / 72)
        pix = page.get_pixmap(matrix=mat, colorspace=fitz.csGRAY)
        img_bytes = pix.tobytes("png")

        try:
            page_text = generate_with_image(img_bytes, OCR_PROMPT)
            texts.append(f"[Page {i + 1}]\n{page_text}")
            print(f"[OCR] Page {i + 1}/{page_count} done ({len(page_text)} chars)")
        except Exception as e:
            print(f"[WARN] OCR failed for page {i + 1}: {e}")

    doc.close()
    return "\n\n".join(texts), page_count


def _extract_docx(filepath: str) -> Tuple[str, int]:
    from docx import Document
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs), max(1, len(paragraphs) // 30)


def _extract_pptx(filepath: str) -> Tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(filepath)
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slide_texts.append("\n".join(texts))
    return "\n\n".join(slide_texts), len(prs.slides)


def _extract_txt(filepath: str) -> Tuple[str, int]:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return text, max(1, len(text) // 3000)


def clean_text(text: str) -> str:
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^\x00-\x7F]+', ' ', text)
    text = re.sub(r'\.{3,}', '...', text)
    text = re.sub(r'-{2,}', '--', text)
    return text.strip()


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    """Split text into overlapping chunks by sentence boundaries."""
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_length = 0

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        if current_length + len(sentence) > chunk_size and current_chunk:
            chunk_text_str = " ".join(current_chunk)
            chunks.append(chunk_text_str)

            # Keep overlap: last few sentences
            overlap_text = " ".join(current_chunk)
            overlap_sentences = []
            overlap_len = 0
            for s in reversed(current_chunk):
                if overlap_len + len(s) <= overlap:
                    overlap_sentences.insert(0, s)
                    overlap_len += len(s)
                else:
                    break
            current_chunk = overlap_sentences
            current_length = sum(len(s) for s in current_chunk)

        current_chunk.append(sentence)
        current_length += len(sentence)

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    # Filter out very short chunks
    chunks = [c for c in chunks if len(c.strip()) > 50]
    return chunks
